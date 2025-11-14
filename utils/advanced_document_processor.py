import os
import logging
import pandas as pd
import PyPDF2
import csv
from io import StringIO, BytesIO
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from pptx import Presentation
import json
import zipfile

logger = logging.getLogger(__name__)

class AdvancedDocumentProcessor:
    """Advanced document processor with OCR and extended format support"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.json': self._process_json,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.csv': self._process_csv,
            '.txt': self._process_text,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.png': self._process_image,
            '.bmp': self._process_image,
            '.tiff': self._process_image
        }

    def process_file(self, file_path, extract_images=False):
        """Process a file and extract text content with advanced features"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            processor = self.supported_formats[file_extension]
            result = processor(file_path, extract_images=extract_images)
            
            # Ensure result is a dictionary with consistent structure
            if isinstance(result, str):
                result = {
                    'text_content': result,
                    'metadata': {},
                    'images': [],
                    'tables': [],
                    'structure': {}
                }
            
            # Clean and validate content
            text_content = result.get('text_content', '')
            if not isinstance(text_content, str):
                text_content = str(text_content)
            
            text_content = self._clean_text(text_content)
            
            if not text_content.strip():
                raise ValueError("No text content could be extracted from the file")
            
            result['text_content'] = text_content
            result['word_count'] = len(text_content.split())
            
            logger.info(f"Successfully processed {file_path}, extracted {len(text_content)} characters")
            return result
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise

    def _process_pdf(self, file_path, extract_images=False):
        """Extract text from PDF file with OCR support"""
        try:
            text_content = []
            metadata = {}
            images = []
            tables = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                metadata['total_pages'] = len(pdf_reader.pages)
                
                if hasattr(pdf_reader, 'metadata') and pdf_reader.metadata:
                    metadata.update({
                        'title': pdf_reader.metadata.get('/Title', ''),
                        'author': pdf_reader.metadata.get('/Author', ''),
                        'creator': pdf_reader.metadata.get('/Creator', ''),
                        'subject': pdf_reader.metadata.get('/Subject', '')
                    })
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    
                    if text.strip():
                        text_content.append(f"Page {page_num + 1}:\n{text}")
                    elif extract_images:
                        # If no text found, try OCR on images
                        # This is a simplified approach - in a real implementation,
                        # you'd extract images from the PDF and process them
                        logger.info(f"No text found on page {page_num + 1}, OCR might be needed")
                        text_content.append(f"Page {page_num + 1}: [Image content - OCR needed]")
            
            result_text = '\n\n'.join(text_content)
            
            return {
                'text_content': result_text,
                'metadata': metadata,
                'images': images,
                'tables': tables,
                'structure': {'pages': metadata['total_pages']}
            }
            
        except Exception as e:
            raise Exception(f"Error processing PDF: {str(e)}")

    def _process_docx(self, file_path, extract_images=False):
        """Extract text from Word documents"""
        try:
            doc = DocxDocument(file_path)
            text_content = []
            tables = []
            metadata = {}
            
            # Extract document properties
            if hasattr(doc, 'core_properties'):
                props = doc.core_properties
                metadata.update({
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'created': str(props.created) if props.created else '',
                    'modified': str(props.modified) if props.modified else ''
                })
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract tables
            for table_num, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    tables.append({
                        'table_number': table_num + 1,
                        'data': table_data,
                        'rows': len(table_data),
                        'columns': len(table_data[0]) if table_data else 0
                    })
                    
                    # Add table content to text
                    text_content.append(f"\nTable {table_num + 1}:")
                    for row in table_data:
                        text_content.append(" | ".join(row))
            
            return {
                'text_content': '\n'.join(text_content),
                'metadata': metadata,
                'images': [],
                'tables': tables,
                'structure': {'paragraphs': len(doc.paragraphs), 'tables': len(tables)}
            }
            
        except Exception as e:
            raise Exception(f"Error processing Word document: {str(e)}")

    def _process_pptx(self, file_path, extract_images=False):
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            text_content = []
            metadata = {}
            slides_data = []
            
            # Extract presentation properties
            if hasattr(prs, 'core_properties'):
                props = prs.core_properties
                metadata.update({
                    'title': props.title or '',
                    'author': props.author or '',
                    'subject': props.subject or '',
                    'created': str(props.created) if props.created else '',
                    'modified': str(props.modified) if props.modified else ''
                })
            
            metadata['total_slides'] = len(prs.slides)
            
            # Extract text from slides
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_title = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not slide_title and hasattr(shape, 'text_frame'):
                            slide_title = shape.text.strip()
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = f"Slide {slide_num + 1}:"
                    if slide_title:
                        slide_content += f" {slide_title}"
                    slide_content += f"\n{chr(10).join(slide_text)}"
                    text_content.append(slide_content)
                    
                    slides_data.append({
                        'slide_number': slide_num + 1,
                        'title': slide_title,
                        'content': slide_text
                    })
            
            return {
                'text_content': '\n\n'.join(text_content),
                'metadata': metadata,
                'images': [],
                'tables': [],
                'structure': {'slides': slides_data}
            }
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint: {str(e)}")

    def _process_image(self, file_path, extract_images=False):
        """Extract text from images using OCR"""
        try:
            # Open and process image
            image = Image.open(file_path)
            
            # Get image metadata
            metadata = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'width': image.width,
                'height': image.height
            }
            
            # Perform OCR
            try:
                text = pytesseract.image_to_string(image, lang='fra+eng+spa')
                confidence = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                
                # Calculate average confidence
                confidences = [int(conf) for conf in confidence['conf'] if int(conf) > 0]
                avg_confidence = sum(confidences) / len(confidences) if confidences else 0
                
                metadata['ocr_confidence'] = avg_confidence
                
            except Exception as ocr_error:
                logger.warning(f"OCR failed for {file_path}: {ocr_error}")
                text = f"[Image file - OCR failed: {str(ocr_error)}]"
                metadata['ocr_confidence'] = 0
            
            return {
                'text_content': text,
                'metadata': metadata,
                'images': [{'path': file_path, 'metadata': metadata}],
                'tables': [],
                'structure': {'image_processed': True}
            }
            
        except Exception as e:
            raise Exception(f"Error processing image: {str(e)}")

    def _process_json(self, file_path, extract_images=False):
        """Extract text from JSON file with enhanced structure analysis"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Analyze JSON structure
            structure = self._analyze_json_structure(data)
            
            # Convert JSON to readable text
            text_content = self._json_to_text(data)
            
            return {
                'text_content': text_content,
                'metadata': {
                    'json_type': type(data).__name__,
                    'total_keys': structure.get('total_keys', 0),
                    'max_depth': structure.get('max_depth', 0)
                },
                'images': [],
                'tables': [],
                'structure': structure
            }
            
        except Exception as e:
            raise Exception(f"Error processing JSON: {str(e)}")

    def _process_excel(self, file_path, extract_images=False):
        """Extract text from Excel file with enhanced table analysis"""
        try:
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            tables = []
            metadata = {'sheet_names': excel_file.sheet_names}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Analyze sheet
                sheet_info = {
                    'name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist(),
                    'has_numeric_data': any(df.dtypes == 'number'),
                    'has_date_data': any(df.dtypes == 'datetime64[ns]')
                }
                
                tables.append(sheet_info)
                
                # Convert DataFrame to text
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += f"Dimensions: {len(df)} rows × {len(df.columns)} columns\n"
                sheet_text += "Columns: " + ", ".join(df.columns.tolist()) + "\n\n"
                
                # Add sample data (first 20 rows)
                sample_size = min(20, len(df))
                for index in range(sample_size):
                    row = df.iloc[index]
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    if row_text.strip():
                        sheet_text += f"Row {index + 1}: {row_text}\n"
                
                if len(df) > sample_size:
                    sheet_text += f"... and {len(df) - sample_size} more rows\n"
                
                text_content.append(sheet_text)
            
            return {
                'text_content': '\n\n'.join(text_content),
                'metadata': metadata,
                'images': [],
                'tables': tables,
                'structure': {'sheets': len(excel_file.sheet_names), 'total_rows': sum(t['rows'] for t in tables)}
            }
            
        except Exception as e:
            raise Exception(f"Error processing Excel: {str(e)}")

    def _process_csv(self, file_path, extract_images=False):
        """Extract text from CSV file with enhanced analysis"""
        try:
            text_content = []
            
            with open(file_path, 'r', encoding='utf-8') as file:
                # Detect delimiter
                sample = file.read(1024)
                file.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.DictReader(file, delimiter=delimiter)
                fieldnames = reader.fieldnames or []
                
                # Analyze CSV structure
                rows_data = list(reader)
                total_rows = len(rows_data)
                
                metadata = {
                    'delimiter': delimiter,
                    'columns': len(fieldnames),
                    'column_names': fieldnames,
                    'total_rows': total_rows
                }
                
                # Add column headers
                text_content.append(f"CSV File Analysis:")
                text_content.append(f"Columns ({len(fieldnames)}): " + ", ".join(fieldnames))
                text_content.append(f"Total rows: {total_rows}")
                text_content.append("")
                
                # Add sample data (first 20 rows)
                sample_size = min(20, total_rows)
                for row_num, row in enumerate(rows_data[:sample_size], 1):
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if val])
                    if row_text.strip():
                        text_content.append(f"Row {row_num}: {row_text}")
                
                if total_rows > sample_size:
                    text_content.append(f"... and {total_rows - sample_size} more rows")
            
            return {
                'text_content': '\n'.join(text_content),
                'metadata': metadata,
                'images': [],
                'tables': [{'name': 'main_table', 'rows': total_rows, 'columns': len(fieldnames)}],
                'structure': {'rows': total_rows, 'columns': len(fieldnames)}
            }
            
        except Exception as e:
            raise Exception(f"Error processing CSV: {str(e)}")

    def _process_text(self, file_path, extract_images=False):
        """Extract text from plain text file with analysis"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
            except Exception as e:
                raise Exception(f"Error reading text file with encoding: {str(e)}")
        
        # Analyze text structure
        lines = content.split('\n')
        words = content.split()
        
        metadata = {
            'total_lines': len(lines),
            'total_words': len(words),
            'total_characters': len(content),
            'encoding_detected': 'utf-8'
        }
        
        return {
            'text_content': content,
            'metadata': metadata,
            'images': [],
            'tables': [],
            'structure': {'lines': len(lines), 'words': len(words)}
        }

    def _analyze_json_structure(self, data, depth=0, max_depth=0):
        """Analyze JSON structure recursively"""
        max_depth = max(max_depth, depth)
        
        if isinstance(data, dict):
            total_keys = len(data)
            for value in data.values():
                if isinstance(value, (dict, list)):
                    child_analysis = self._analyze_json_structure(value, depth + 1, max_depth)
                    total_keys += child_analysis.get('total_keys', 0)
                    max_depth = max(max_depth, child_analysis.get('max_depth', 0))
        elif isinstance(data, list):
            total_keys = len(data)
            for item in data:
                if isinstance(item, (dict, list)):
                    child_analysis = self._analyze_json_structure(item, depth + 1, max_depth)
                    total_keys += child_analysis.get('total_keys', 0)
                    max_depth = max(max_depth, child_analysis.get('max_depth', 0))
        else:
            total_keys = 0
        
        return {'total_keys': total_keys, 'max_depth': max_depth}

    def _json_to_text(self, data, level=0):
        """Convert JSON data to readable text with better formatting"""
        text_content = []
        indent = "  " * level
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    text_content.append(f"{indent}{key}:")
                    text_content.append(self._json_to_text(value, level + 1))
                else:
                    text_content.append(f"{indent}{key}: {value}")
        
        elif isinstance(data, list):
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    text_content.append(f"{indent}[{i}]:")
                    text_content.append(self._json_to_text(item, level + 1))
                else:
                    text_content.append(f"{indent}[{i}]: {item}")
        
        else:
            text_content.append(f"{indent}{data}")
        
        return '\n'.join(text_content)

    def _clean_text(self, text):
        """Clean and normalize text content"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            cleaned_line = line.strip()
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
        
        # Join lines and limit excessive newlines
        cleaned_text = '\n'.join(cleaned_lines)
        
        # Remove excessive consecutive newlines
        while '\n\n\n' in cleaned_text:
            cleaned_text = cleaned_text.replace('\n\n\n', '\n\n')
        
        return cleaned_text

    def extract_tables_from_text(self, text):
        """Extract table-like structures from text"""
        tables = []
        lines = text.split('\n')
        
        current_table = []
        in_table = False
        
        for line in lines:
            # Simple heuristic: if line contains multiple | or tabs, it might be a table row
            if '|' in line and line.count('|') >= 2:
                current_table.append(line.split('|'))
                in_table = True
            elif '\t' in line and line.count('\t') >= 2:
                current_table.append(line.split('\t'))
                in_table = True
            else:
                if in_table and len(current_table) > 1:
                    # Save completed table
                    tables.append({
                        'rows': len(current_table),
                        'columns': len(current_table[0]) if current_table else 0,
                        'data': current_table
                    })
                current_table = []
                in_table = False
        
        # Don't forget the last table if file ends with table
        if in_table and len(current_table) > 1:
            tables.append({
                'rows': len(current_table),
                'columns': len(current_table[0]) if current_table else 0,
                'data': current_table
            })
        
        return tables